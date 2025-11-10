import os, io, uuid, streamlit as st
from dotenv import load_dotenv
from pinecone import Pinecone, ServerlessSpec
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from pptx import Presentation
import google.generativeai as genai

# ---------- setup ----------
load_dotenv()

# Get API keys from .env (local) or Streamlit secrets (cloud)
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY") or st.secrets.get("PINECONE_API_KEY")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY")
INDEX_NAME = os.getenv("INDEX_NAME") or st.secrets.get("INDEX_NAME", "jarvis-index")

genai.configure(api_key=GEMINI_API_KEY)

EMBED_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
EMBED_DIM = 384

@st.cache_resource
def get_embedder():
    return SentenceTransformer(EMBED_MODEL_NAME)

@st.cache_resource
def get_pinecone_index():
    pc = Pinecone(api_key=PINECONE_API_KEY)
    names = [ix.name for ix in pc.list_indexes()]
    if INDEX_NAME not in names:
        pc.create_index(
            name=INDEX_NAME,
            dimension=EMBED_DIM,
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1"),
        )
    return pc.Index(INDEX_NAME)

embedder = get_embedder()
index = get_pinecone_index()

# ---------- helper functions ----------
def chunk_text(text, size=800, overlap=150):
    """Splits text into chunks more safely for low-memory systems."""
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + size, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start += size - overlap
    return chunks

def read_pdf_bytes(file_bytes):
    reader = PdfReader(io.BytesIO(file_bytes))
    return "\n".join([p.extract_text() or "" for p in reader.pages])

def read_pptx_bytes(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        buf = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        if buf:
            parts.append(f"[Slide {i}] " + "\n".join(buf))
    return "\n".join(parts)

def upsert_texts(chunks, src):
    batch = []
    for chunk in chunks:
        emb = embedder.encode([chunk], normalize_embeddings=True)[0].tolist()
        batch.append({
            "id": str(uuid.uuid4()),
            "values": emb,
            "metadata": {"text": chunk, "source": src}
        })
        if len(batch) >= 200:
            index.upsert(vectors=batch)
            batch = []
    if batch:
        index.upsert(vectors=batch)

def retrieve(q, top_k=5):
    q_emb = embedder.encode([q], normalize_embeddings=True)[0].tolist()
    res = index.query(vector=q_emb, top_k=top_k, include_metadata=True)
    snippets = [m["metadata"]["text"] for m in res["matches"]]
    return snippets

# ---------- Gemini response ----------
def ask_gemini(question, context_snippets):
    import google.generativeai as genai
    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

    context = "\n\n".join(f"- {c}" for c in context_snippets)
    prompt = f"""Use the context below to answer the question. Provide a comprehensive answer based on the context provided.
If the specific information isn't in the context, you can say: "This specific information isn't covered in the uploaded notes, but based on the context provided..."

Context:
{context}

Question: {question}
Answer:"""

    # ✅ Gemini Flash (fast & free)
    try:
        model = genai.GenerativeModel("models/gemini-flash-latest")
        response = model.generate_content(
            prompt,
            generation_config={
                "temperature": 0.7,
                "max_output_tokens": 4096,
            }
        )
        
        # Check if response has valid content
        if response.candidates and len(response.candidates) > 0:
            candidate = response.candidates[0]
            if candidate.content and candidate.content.parts:
                return candidate.content.parts[0].text.strip()
        
        # If no valid content, provide helpful message
        return "I couldn't generate a response. This might be due to content safety filters or insufficient context."
    except Exception as e:
        return f"⚠️ Error: {str(e)}"


# ---------- Streamlit UI ----------
st.set_page_config(page_title="Jarvis (Gemini Model)", page_icon="🤖")
tab1, tab2 = st.tabs(["📤 Upload", "💬 Chat"])

with tab1:
    st.header("Upload your notes (PDF or PPTX)")
    files = st.file_uploader("Upload files", type=["pdf", "pptx"], accept_multiple_files=True)
    if st.button("Send to Pinecone") and files:
        with st.spinner("Processing and uploading..."):
            total = 0
            for f in files:
                data = f.read()
                text = read_pdf_bytes(data) if f.name.endswith(".pdf") else read_pptx_bytes(data)
                chunks = chunk_text(text)
                upsert_texts(chunks, f.name)
                total += len(chunks)
        st.success(f"Uploaded {len(files)} file(s) ({total} chunks).")

with tab2:
    st.header("Chat with your notes")
    
    # Initialize session state for chat history
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    
    # Display chat history
    for i, (question, answer) in enumerate(st.session_state.chat_history):
        st.write(f"**You:** {question}")
        st.write(f"**🤖 Jarvis (Gemini):** {answer}")
        st.divider()
    
    # Input section for new question
    col1, col2 = st.columns([4, 1])
    with col1:
        q = st.text_input("Ask your question:", key=f"input_{len(st.session_state.chat_history)}")
    with col2:
        ask_button = st.button("Ask", key=f"button_{len(st.session_state.chat_history)}")
    
    if ask_button and q:
        with st.spinner("Thinking..."):
            ctx = retrieve(q)
            if not ctx:
                st.warning("No data found. Please upload files first.")
            else:
                ans = ask_gemini(q, ctx)
                # Add to chat history
                st.session_state.chat_history.append((q, ans))
                # Rerun to display new message and clear input
                st.rerun()
